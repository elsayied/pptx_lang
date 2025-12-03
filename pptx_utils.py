import asyncio
import base64
import json
import logging
import os
import re
import zlib
from pprint import pprint
from typing import Dict, List

import httpx
import nest_asyncio
from elevenlabs.client import ElevenLabs
from google import genai
from pptx import Presentation
from pptx.util import Inches, Pt
from youtube_transcript_api import YouTubeTranscriptApi

# Apply nest_asyncio to allow nested event loops
nest_asyncio.apply()

# --- API KEYS (SECURITY WARNING: DO NOT USE IN PUBLIC REPOSITORIES) ---


# ---

logging.basicConfig(level=logging.INFO)

# --- Docling Availability Check ---
try:
    from docling.document_converter import DocumentConverter

    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False
    logging.warning("Docling not available. Mock content may be used.")

# --- PROMPTS ---
GEMINI_LECTURE_PROMPT = """
You are a master presentation creator. Transform the following text into a presentation in Markdown.
Output ONLY the Markdown.

Syntax:
1. Slides separated by `---`
2. Title prefixed with `#`
3. `layout: title_content` (or others) on first line.
4. `::: notes` for speaker notes.
5. `::: column` for columns.
6. **Diagrams with Mermaid**: To add a diagram, use a fenced code block with the `mermaid` language identifier. Create diagrams that are clear, well-structured, and visually impressive.
    *   **Use different shapes**: `A[Client]`, `B(Database)`, `C{Decision}`.
    *   **Add styles**: Use `classDef` to define styles for nodes and `class` to apply them.
    *   **Choose the right layout**: `TD` (top-down), `LR` (left-right), etc.

    Example of a beautiful Mermaid diagram:
    ```mermaid
    graph TD;
        subgraph "User Interaction"
            A[Start] --> B{User Login?};
            B -- Yes --> C[Access Dashboard];
            B -- No --> D[Show Login Page];
        end

        subgraph "Backend Services"
            C --> E(API Gateway);
            E --> F[Auth Service];
            E --> G[Data Service];
        end

        classDef start-end fill:#2E8B57,stroke:#333,stroke-width:2px,color:#fff;
        classDef process fill:#4682B4,stroke:#333,stroke-width:2px,color:#fff;
        classDef decision fill:#DAA520,stroke:#333,stroke-width:2px,color:#fff;

        class A,C start-end;
        class B decision;
        class D,E,F,G process;
    ```

Text to transform:
"""

GEMINI_PODCAST_PROMPT = """
You are a podcast producer. Turn the text into a dialogue script between 'Sascha' and 'Marina'.
Output JSON ONLY: [{"speaker": "Sascha", "text": "..."}]

Text:
{text_content}
"""

# --- HELPER FUNCTIONS ---


def create_table_from_markdown(text: str) -> List[List[str]]:
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    if not lines:
        return []
    lines = [l for l in lines if not re.match(r"^\s*\|?.*--.*\|?\s*$", l)]
    table_data = []
    for row_str in lines:
        if row_str.startswith("|"):
            row_str = row_str[1:]
        if row_str.endswith("|"):
            row_str = row_str[:-1]
        cells = [cell.strip() for cell in row_str.split("|")]
        table_data.append(cells)
    return table_data


def add_formatted_text_runs(paragraph, text, bold=False, italic=False, underline=False):
    parts = re.split(
        r"(\*\*\*[\s\S]+?\*\*\*|\*\*[\s\S]+?\*\*|\*[\s\S]+?\*|__[\s\S]+?__)", text
    )
    if len(parts) == 1:
        if text:
            run = paragraph.add_run()
            run.text = text
            font = run.font
            font.bold = bold
            font.italic = italic
            font.underline = underline
        return
    for part in parts:
        if not part:
            continue
        if part.startswith("***") and part.endswith("***"):
            add_formatted_text_runs(
                paragraph, part[3:-3], bold=True, italic=True, underline=underline
            )
        elif part.startswith("**") and part.endswith("**"):
            add_formatted_text_runs(
                paragraph, part[2:-2], bold=True, italic=italic, underline=underline
            )
        elif part.startswith("*") and part.endswith("*"):
            add_formatted_text_runs(
                paragraph, part[1:-1], bold=bold, italic=True, underline=underline
            )
        elif part.startswith("__") and part.endswith("__"):
            add_formatted_text_runs(
                paragraph, part[2:-2], bold=bold, italic=italic, underline=True
            )
        else:
            add_formatted_text_runs(
                paragraph, part, bold=bold, italic=italic, underline=underline
            )


def add_bullet_points_from_markdown(text_frame, points: str):
    if not text_frame.text.strip():
        text_frame.text = ""

    def get_level_and_text(line: str) -> tuple[int, str]:
        stripped_line = line.lstrip()
        text = stripped_line
        if text.startswith(("-", "*", "+")) and text[1:2] in (" ", ""):
            text = text[1:].lstrip()
        indent = len(line) - len(line.lstrip())
        return indent // 2, text

    lines = [line for line in points.split("\n") if line.strip()]
    for line in lines:
        level, text = get_level_and_text(line)
        p = text_frame.add_paragraph()
        add_formatted_text_runs(p, text)
        p.level = min(level, 8)


def parse_markdown_to_slides(content: str) -> List[Dict]:
    slides = []
    slide_contents = re.split(r"\n---\n", content)
    for slide_content in slide_contents:
        if not slide_content.strip():
            continue
        lines = slide_content.strip().split("\n")
        slide_data = {
            "layout": "title_content",
            "title": None,
            "blocks": [],
            "columns": [],
            "notes": "",
        }
        if lines and lines[0].startswith("layout:"):
            slide_data["layout"] = lines[0].split(":", 1)[1].strip()
            lines.pop(0)
        if lines and lines[0].startswith("# "):
            slide_data["title"] = lines[0][2:].strip()
            lines.pop(0)

        line_idx = 0
        in_notes = False
        in_col = False
        while line_idx < len(lines):
            line = lines[line_idx]
            stripped = line.strip()
            if stripped.lower().startswith("::: notes"):
                in_notes = True
                line_idx += 1
                continue
            if stripped.lower().startswith("::: column"):
                in_col = True
                if not slide_data["columns"] and slide_data["blocks"]:
                    slide_data["columns"].append(slide_data["blocks"])
                    slide_data["blocks"] = []
                slide_data["columns"].append([])
                line_idx += 1
                continue
            if stripped == ":::":
                in_notes = False
                in_col = False
                line_idx += 1
                continue
            if in_notes:
                slide_data["notes"] += line + "\n"
                line_idx += 1
                continue

            curr = (
                slide_data["columns"][-1]
                if in_col and slide_data["columns"]
                else slide_data["blocks"]
            )
            if not stripped:
                line_idx += 1
                continue

            if stripped.startswith("```"):
                code = []
                lang = stripped[3:]
                line_idx += 1
                while line_idx < len(lines) and not lines[line_idx].strip().startswith(
                    "```"
                ):
                    code.append(lines[line_idx])
                    line_idx += 1
                curr.append(
                    {"type": "code", "language": lang, "content": "\n".join(code)}
                )
                line_idx += 1
                continue

            if re.match(r"^\s*!\[.*\]\(.*\)\s*$", line):
                curr.append({"type": "image", "content": line.strip()})
                line_idx += 1
                continue

            if line.lstrip().startswith(("-", "*", "+")):
                bullets = []
                start_indent = len(line) - len(line.lstrip())
                while line_idx < len(lines):
                    curr_l = lines[line_idx]
                    curr_s = curr_l.strip()
                    if curr_s.startswith((":::", "```")):
                        break
                    if curr_s and (len(curr_l) - len(curr_l.lstrip()) < start_indent):
                        break
                    bullets.append(curr_l)
                    line_idx += 1
                curr.append({"type": "bullet", "content": "\n".join(bullets)})
                continue

            text_lines = []
            while line_idx < len(lines):
                curr_l = lines[line_idx]
                curr_s = curr_l.strip()
                if (
                    not curr_s
                    or curr_l.lstrip().startswith(("-", "*", "+"))
                    or re.match(r"^\s*!\[.*\]\(.*\)\s*$", curr_l)
                    or curr_s.startswith(("```", ":::"))
                ):
                    break
                text_lines.append(curr_l)
                line_idx += 1
            if text_lines:
                curr.append({"type": "text", "content": "\n".join(text_lines)})
        slides.append(slide_data)
    return slides


def render_mermaid_diagram(mermaid_code: str, output_path: str):
    """Renders a Mermaid diagram using the Kroki.io service."""
    try:
        encoded_diagram = base64.urlsafe_b64encode(
            zlib.compress(mermaid_code.encode("utf-8"), 9)
        ).decode("utf-8")

        kroki_url = f"https://kroki.io/mermaid/png/{encoded_diagram}"

        with httpx.Client() as client:
            response = client.get(kroki_url)
            response.raise_for_status()

            with open(output_path, "wb") as f:
                f.write(response.content)

        logging.info(
            f"Successfully rendered Mermaid diagram to {output_path} using Kroki"
        )
        return True

    except Exception as e:
        logging.error(f"Failed to render Mermaid diagram via Kroki: {e}")
        try:
            from PIL import Image, ImageDraw

            img = Image.new("RGB", (600, 300), color=(255, 255, 255))
            d = ImageDraw.Draw(img)
            d.text((10, 10), f"Mermaid/Kroki Render Error:\n{e}", fill=(255, 0, 0))
            img.save(output_path)
        except Exception as img_e:
            logging.error(f"Failed to create placeholder image: {img_e}")
        return False


def preprocess_markdown_for_diagrams(markdown_content: str) -> str:
    """
    Finds Mermaid code blocks, renders them as images, and replaces
    the code block with a Markdown image tag.
    """
    pattern = r"```mermaid\n([\s\S]+?)\n```"

    mermaid_codes = re.findall(pattern, markdown_content)
    if not mermaid_codes:
        return markdown_content

    images_dir = "generated_images"
    os.makedirs(images_dir, exist_ok=True)

    modified_content = markdown_content

    for i, mermaid_code in enumerate(mermaid_codes):
        image_filename = f"generated_diagram_{i}.png"
        output_path = os.path.join(images_dir, image_filename)

        render_mermaid_diagram(mermaid_code, output_path)

        original_block = f"```mermaid\n{mermaid_code}\n```"
        new_tag = f"![Mermaid Diagram]({output_path})"
        modified_content = modified_content.replace(original_block, new_tag, 1)

    return modified_content


def create_presentation_from_markdown(
    content: str, output_path: str = "output.pptx"
) -> str:
    content = preprocess_markdown_for_diagrams(content)
    prs = Presentation()
    layout_map = {
        "title_slide": prs.slide_layouts[0],
        "title_content": prs.slide_layouts[1],
        "section_header": prs.slide_layouts[2],
        "two_content": prs.slide_layouts[3],
        "comparison": prs.slide_layouts[3],
        "title_only": prs.slide_layouts[5],
        "blank": prs.slide_layouts[6],
        "picture_and_caption": prs.slide_layouts[8],
    }
    slides_data = parse_markdown_to_slides(content)

    for slide_data in slides_data:
        layout = layout_map.get(slide_data["layout"], layout_map["title_content"])
        slide = prs.slides.add_slide(layout)
        if slide_data["title"] and slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]
        if slide_data["notes"]:
            slide.notes_slide.notes_text_frame.text = slide_data["notes"]

        body = next(
            (
                s
                for s in slide.placeholders
                if s.placeholder_format.idx != 0 and s.has_text_frame
            ),
            None,
        )
        blocks = [b for b in slide_data["blocks"] if b["type"] != "image"]
        if body and blocks:
            tf = body.text_frame
            tf.clear()
            for b in blocks:
                if b["type"] == "text":
                    add_formatted_text_runs(tf.add_paragraph(), b["content"])
                elif b["type"] == "bullet":
                    add_bullet_points_from_markdown(tf, b["content"])
                elif b["type"] == "code":
                    p = tf.add_paragraph()
                    r = p.add_run()
                    r.text = b["content"]
                    r.font.name = "Courier New"
                    r.font.size = Pt(10)

    prs.save(output_path)
    return output_path


# --- CORE FUNCTIONS ---


def generate_structured_markdown(text: str, api_key: str) -> str:
    """Generates slides markdown using Gemini."""
    if not api_key:
        return "Error: Gemini API key not provided."
    client = genai.Client(api_key=api_key)
    try:
        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=GEMINI_LECTURE_PROMPT + text,
        )
        return response.text
    except Exception as e:
        return f"Gemini API Error: {e}"


def generate_podcast_script(raw_text: str, api_key: str, prompt: str = None) -> List[Dict]:
    """Generates script (1 API Call) with robust cleanup."""
    # if not api_key:
    #     logging.error("Gemini API key not provided for podcast script generation.")
    #     return [{"speaker": "Error", "text": "Gemini API key not provided."}]
    # # the error is here
    client = genai.Client(api_key=api_key)

    try:
        if prompt:
            final_prompt = prompt.format(text_content=raw_text)
        else:
            final_prompt = GEMINI_PODCAST_PROMPT.format(text_content=raw_text)

        response = client.models.generate_content(
            model="gemini-2.5-flash",
            contents=final_prompt,
            config={"response_mime_type": "application/json"},
        )
        logging.info(f"Gemini Podcast Script Response: {response.text}")

        text_resp = response.text.strip()
        if text_resp.startswith("```"):
            text_resp = re.sub(
                r"^```(json)?|```$", "", text_resp, flags=re.MULTILINE
            ).strip()

        script = json.loads(text_resp)

        if isinstance(script, dict):
            script = [script]

        return script
    except Exception as e:
        logging.error(f"Script Error: {e}")
        return [
            {
                "speaker": "Error",
                "text": f"Could not generate podcast script. Error: {e}",
            }
        ]


def generate_audio_overview(
    script: List[Dict], output_path: str, api_key: str, voice_mapping: Dict[str, str]
):
    """
    Generates an audio overview of the script using the Eleven Labs API.
    voice_mapping should contain speaker names mapped to their voice_ids.
    Returns True on success, "credit_error" on credit issues, and False on other failures.
    """
    if not api_key:
        logging.error("Eleven Labs API key is not set.")
        return False

    try:
        client = ElevenLabs(api_key=api_key)

        default_voice_id = next(
            (vid for speaker, vid in voice_mapping.items() if speaker == "Marina"), None
        )
        if not default_voice_id:
            all_voices = client.voices.get_all().voices
            default_voice_id = next(
                (v.voice_id for v in all_voices if v.name == "Bella"),
                all_voices[0].voice_id,
            )

        full_audio = b""
        for line in script:
            text = line.get("text", "")
            speaker = line.get("speaker")
            if not text or not speaker:
                continue

            voice_id = voice_mapping.get(speaker, default_voice_id)

            audio_stream = client.text_to_speech.convert(
                text=text, voice_id=voice_id, model_id="eleven_multilingual_v2"
            )

            for chunk in audio_stream:
                full_audio += chunk

        if not full_audio:
            logging.warning("No audio was generated by Eleven Labs.")
            return False

    except Exception as e:
        logging.error(f"Error during Eleven Labs audio generation: {e}")
        if "credit" in str(e).lower():
            return "credit_error"
        return False

    with open(output_path, "wb") as f:
        f.write(full_audio)
    return True


def extract_content_with_docling(
    file_path: str, enabled_ocr=True, page_range: str = None
):
    """
    Extracts content from various file types using the docling library.
    """
    suffix = file_path.split(".")[-1]
    try:
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import (
            PdfPipelineOptions,
            TesseractCliOcrOptions,
        )
        from docling.document_converter import DocumentConverter, PdfFormatOption

        ocr_options = TesseractCliOcrOptions(lang=["eng"])
        pipeline_options = PdfPipelineOptions(
            do_ocr=enabled_ocr, do_table_structure=True, ocr_options=ocr_options
        )
        doc_converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(
                    pipeline_options=pipeline_options,
                )
            }
        )
        doc = doc_converter.convert(file_path).document

        if page_range:
            pages = []
            for part in page_range.split(","):
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    pages.extend(range(start - 1, end))
                else:
                    pages.append(int(part) - 1)

            filtered_content = []
            for i, page in enumerate(doc.pages):
                if i in pages:
                    filtered_content.append(page.export_to_markdown())
            return "\n".join(filtered_content)
        else:
            return doc.export_to_markdown()

    except Exception as e:
        print(f"Error extracting content with docling from {file_path}: {e}")
        return {"error": str(e)}


def time_to_seconds(time_str):
    """Converts a time string of the format M:S or H:M:S to seconds."""
    if not time_str:
        return 0
    parts = time_str.split(":")
    if len(parts) == 2:
        return int(parts[0]) * 60 + int(parts[1])
    elif len(parts) == 3:
        return int(parts[0]) * 3600 + int(parts[1]) * 60 + int(parts[2])
    return 0


def extract_content_from_youtube(url: str, start_time: str, end_time: str):
    """
    Fetches the transcript of a YouTube video and extracts a specific time range.
    """
    try:
        video_id = url.split("v=")[1]
        transcript = YouTubeTranscriptApi.get_transcript(video_id)

        start_seconds = time_to_seconds(start_time)
        end_seconds = time_to_seconds(end_time)

        content = ""
        for item in transcript:
            item_start_time = item["start"]
            if end_seconds:
                if item_start_time >= start_seconds and item_start_time <= end_seconds:
                    content += item["text"] + " "
            elif item_start_time >= start_seconds:
                content += item["text"] + " "

        return content

    except Exception as e:
        logging.error(f"Error processing YouTube video: {e}")
        return f"Error: {e}"
