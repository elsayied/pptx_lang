import asyncio
import json
import logging
import os
import re
from pprint import pprint
from typing import Dict, List

import nest_asyncio
from google import genai
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import asyncio
from mermaid.py_mermaid import MermaidAPI

# Apply nest_asyncio to allow nested event loops (crucial for Streamlit + edge-tts)
nest_asyncio.apply()

# --- HARDCODED API KEY (SECURITY WARNING: DO NOT USE IN PUBLIC REPOSITORIES) ---
# ⚠️ REPLACE THE PLACEHOLDER BELOW WITH YOUR ACTUAL GEMINI API KEY ⚠️
HARDCODED_GEMINI_API_KEY = "AIzaSyA4YsTnbNjl2gKn20EqPa-9nom9yymEwd0"
# ---

logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)
# --- Docling Availability Check ---
try:
    from docling.document_converter import DocumentConverter

    DOCLING_AVAILABLE = True
except ImportError:
    DOCLING_AVAILABLE = False
    logging.warning("Docling not available. Mock content may be used.")

# --- PROMPTS ---

GEMINI_LECTURE_PROMPT = """
You are a master presentation creator, an expert in distilling complex information into clear, engaging, and visually appealing slides. Your mission is to transform the following text into a presentation in a specialized Markdown format.

The presentation should be professional yet captivating. Use concise language, structure information logically, and think about the visual presentation on a slide.

Your output MUST be only the Markdown text and nothing else.

Here is the extended Markdown syntax you will use:

1.  **Slides**: Each slide is separated by a line containing only `---`.

2.  **Title**: A slide's title is prefixed with `#`. Example: `# A Compelling Title`.

3.  **Layouts**: Specify a slide layout with `layout: <layout_name>` on the first line. Supported layouts are: `title_slide`, `title_content` (default), `section_header`, `two_content`, `comparison`, `title_only`, `blank`, `picture_and_caption`.

4.  **Diagrams with Mermaid**: To add a diagram, use a fenced code block with the `mermaid` language identifier. Create diagrams that are clear, well-structured, and visually impressive.
    *   **Use different shapes**: `A[Client]`, `B(Database)`, `C{Decision}`.
    *   **Add styles**: Use `classDef` to define styles for nodes and `class` to apply them.
    *   **Choose the right layout**: `TD` (top-down), `LR` (left-right), etc.

    Example of a beautiful Mermaid diagram:
    ```mermaid
    graph TD;
        subgraph "User Interaction"
            A[Start] --> B{User Login?};
            B -- Yes --> C[Access Dashboard];
            B -- No --> D[Show Login Page];
        end

        subgraph "Backend Services"
            C --> E(API Gateway);
            E --> F[Auth Service];
            E --> G[Data Service];
        end

        classDef start-end fill:#2E8B57,stroke:#333,stroke-width:2px,color:#fff;
        classDef process fill:#4682B4,stroke:#333,stroke-width:2px,color:#fff;
        classDef decision fill:#DAA520,stroke:#333,stroke-width:2px,color:#fff;

        class A,C start-end;
        class B decision;
        class D,E,F,G process;
    ```

5.  **Speaker Notes**: Add non-visible notes for the presenter inside a `::: notes` block.

6.  **Multi-Column Layouts**: Use `::: column` blocks for complex multi-column content.

Now, take the following text and create a brilliant presentation from it:
"""
GEMINI_PODCAST_PROMPT = """
You are a podcast producer. Turn the text into a dialogue script between 'Sascha' and 'Marina'.
Output JSON ONLY: [{"speaker": "Sascha", "text": "..."}]

Text:
{text_content}
"""


def generate_image(prompt: str, output_path: str):
    """Generates an image using Gemini and saves it to a file."""
    try:
        logging.info(f"Generating image for prompt: {prompt}")
        client = genai.Client(api_key=HARDCODED_GEMINI_API_KEY)
        response = client.models.generate_content(
            model="gemini-2.0-flash-preview-image-generation",  # Or another model, as per docs
            contents=[prompt],
        )

        for part in response.parts:
            if part.inline_data is not None:
                image = part.as_image()
                image.save(output_path)
                logging.info(f"Saved generated image to {output_path}")
                return output_path
    except Exception as e:
        logging.error(f"Error generating image: {e}")
        return None


def preprocess_markdown_for_images(markdown_content: str) -> str:
    """
    Finds image generation directives, generates images, and replaces
    the directive with the path to the generated image file.api_key="AIzaSyA4YsTnbNjl2gKn20EqPa-9nom9yymEwd0"
    """
    # Regex to find ![alt text](gemini:prompt)
    pattern = r"!\[(.*?)\]\(gemini:(.*?)\)"
    matches = re.findall(pattern, markdown_content)

    if not matches:
        return markdown_content

    images_dir = "generated_images"
    os.makedirs(images_dir, exist_ok=True)

    modified_content = markdown_content
    for i, (alt_text, prompt) in enumerate(matches):
        image_filename = f"generated_image_{i}.png"
        output_path = os.path.join(images_dir, image_filename)
        sleep_min()
        if generate_image(prompt, output_path):
            # Replace the gemini directive with the local file path
            original_directive = f"![{alt_text}](gemini:{prompt})"
            new_directive = f"![{alt_text}]({output_path})"
            modified_content = modified_content.replace(
                original_directive, new_directive, 1
            )
            logging.info(f"{new_directive = }\n {original_directive = }\n")
        else:
            logging.warning(f"Could not generate image for prompt: {prompt}")

    return modified_content


def create_table_from_markdown(text: str) -> List[List[str]]:
    """Convert Markdown table to table data."""
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    if not lines:
        return []

    # Filter out separator line
    lines = [l for l in lines if not re.match(r"^\s*\|?.*--.*\|?\s*$", l)]

    table_data = []
    for row_str in lines:
        if row_str.startswith("|"):
            row_str = row_str[1:]
        if row_str.endswith("|"):
            row_str = row_str[:-1]

        cells = [cell.strip() for cell in row_str.split("|")]
        table_data.append(cells)

    return table_data


def add_formatted_text_runs(paragraph, text, bold=False, italic=False, underline=False):
    """
    Recursively parses markdown-like formatting for bold, italic, and
    underline, and adds formatted runs to the given paragraph.

    This function works by repeatedly splitting the text by all possible
    formatting markers. If the text can't be split further, it's considered
    plain text and added as a single run with the formatting inherited from
    its parent calls. If the text is split, the function calls itself for
    each segment, updating the formatting flags (bold, italic, etc.) for
    segments that were wrapped in markers.

    - ***text*** for bold & italic
    - **text** for bold
    - *text* for italic
    - __text__ for underline
    """
    # Split by the formatting markers, keeping them.
    parts = re.split(
        r"(\*\*\*[\s\S]+?\*\*\*|\*\*[\s\S]+?\*\*|\*[\s\S]+?\*|__[\s\S]+?__)", text
    )

    # Base case: If the text is not split into more than one part,
    # it means no formatting markers were found. Add the text as a single run.
    if len(parts) == 1:
        if text:
            run = paragraph.add_run()
            run.text = text
            font = run.font
            font.bold = bold
            font.italic = italic
            font.underline = underline
        return

    # Recursive step: For each part of the split text, determine if it's
    # a formatted segment or plain text, and recurse.
    for part in parts:
        if not part:
            continue

        # Order of checks is important: *** must be checked before ** or *.
        if part.startswith("***") and part.endswith("***"):
            # It's a bold and italic segment, recurse on its content with new formatting.
            add_formatted_text_runs(
                paragraph, part[3:-3], bold=True, italic=True, underline=underline
            )
        elif part.startswith("**") and part.endswith("**"):
            # It's a bold segment, recurse on its content.
            add_formatted_text_runs(
                paragraph, part[2:-2], bold=True, italic=italic, underline=underline
            )
        elif part.startswith("*") and part.endswith("*"):
            # It's an italic segment, recurse on its content.
            add_formatted_text_runs(
                paragraph, part[1:-1], bold=bold, italic=True, underline=underline
            )
        elif part.startswith("__") and part.endswith("__"):
            # It's an underline segment, recurse on its content.
            add_formatted_text_runs(
                paragraph, part[2:-2], bold=bold, italic=italic, underline=True
            )
        else:
            # This is a plain text part with no new markers.
            # Recurse on it, passing down the formatting from the parent.
            add_formatted_text_runs(
                paragraph, part, bold=bold, italic=italic, underline=underline
            )


def add_bullet_points_from_markdown(text_frame, points: str):
    """Add bullet points to a text frame from Markdown list."""
    if not text_frame.text.strip():
        text_frame.text = ""

    def get_level_and_text(line: str) -> tuple[int, str]:
        """Determine level and clean text from a Markdown list item."""
        stripped_line = line.lstrip()
        text = stripped_line

        if text.startswith(("-", "*", "+")) and text[1:2] in (" ", ""):
            text = text[1:].lstrip()

        indent = len(line) - len(line.lstrip())
        level = indent // 2
        return level, text

    lines = [line for line in points.split("\n") if line.strip()]
    if not lines:
        return

    for line in lines:
        level, text = get_level_and_text(line)
        p = text_frame.add_paragraph()
        add_formatted_text_runs(p, text)
        p.level = min(level, 8)


def parse_markdown_to_slides(content: str) -> List[Dict]:
    """
    Parses markdown-like text into a list of slide definitions.

    This function iterates through the input content line by line, maintaining
    state (e.g., `in_notes_block`, `in_column_block`) to determine what kind of
    content is currently being parsed. It identifies structural elements like
    layout definitions, titles, and various content blocks (notes, columns,
    code, lists, tables, text), converting them into a structured list of
    dictionaries, where each dictionary represents a single slide.
    """
    slides = []
    # A slide is defined as the content between '---' separators.
    import re


from typing import Dict, List


def parse_markdown_to_slides(content: str) -> List[Dict]:
    slides = []
    slide_contents = re.split(r"\n---\n", content)

    for slide_content in slide_contents:
        if not slide_content.strip():
            continue

        lines = slide_content.strip().split("\n")

        slide_data = {
            "layout": "title_content",
            "title": None,
            "blocks": [],
            "columns": [],
            "notes": "",
        }

        # The first lines of a slide can define the layout and title.
        if lines and lines[0].startswith("layout:"):
            slide_data["layout"] = lines[0].split(":", 1)[1].strip()
            lines.pop(0)

        if lines and lines[0].startswith("# "):
            slide_data["title"] = lines[0][2:].strip()
            lines.pop(0)

        line_idx = 0
        in_notes_block = False
        in_column_block = False

        while line_idx < len(lines):
            line = lines[line_idx]
            stripped_line = line.strip()  # Calculate this once

            # 1. ROBUST MARKER CHECKS (Handle extra spaces/case)
            if stripped_line.lower().startswith("::: notes"):
                in_notes_block = True
                line_idx += 1
                continue

            if stripped_line.lower().startswith("::: column"):
                in_column_block = True
                if not slide_data["columns"] and slide_data["blocks"]:
                    slide_data["columns"].append(slide_data["blocks"])
                    slide_data["blocks"] = []
                slide_data["columns"].append([])
                line_idx += 1
                continue

            if stripped_line == ":::":
                in_notes_block = False
                in_column_block = False
                line_idx += 1
                continue

            if in_notes_block:
                slide_data["notes"] += line + "\n"
                line_idx += 1
                continue

            current_blocks = (
                slide_data["columns"][-1]
                if in_column_block and slide_data["columns"]
                else slide_data["blocks"]
            )

            if not stripped_line:
                line_idx += 1
                continue

            # --- PARSING BLOCKS ---

            # Code blocks
            if stripped_line.startswith("```"):
                code_lines = []
                lang = stripped_line[3:]
                line_idx += 1
                while line_idx < len(lines) and not lines[line_idx].strip().startswith(
                    "```"
                ):
                    code_lines.append(lines[line_idx])
                    line_idx += 1
                line_idx += 1
                current_blocks.append(
                    {"type": "code", "language": lang, "content": "\n".join(code_lines)}
                )
                continue

            # Image
            if re.match(r"^\s*!\[.*\]\(.*\)\s*$", line):
                current_blocks.append({"type": "image", "content": line.strip()})
                line_idx += 1
                continue

            # Bulleted lists
            if line.lstrip().startswith(("-", "*", "+")):
                bullet_lines = []
                start_indent = len(line) - len(line.lstrip())

                while line_idx < len(lines):
                    curr_line = lines[line_idx]
                    curr_stripped = curr_line.strip()

                    # --- CRITICAL FIX 1: Stop if we hit a marker ---
                    if curr_stripped.startswith(":::") or curr_stripped.startswith(
                        "```"
                    ):
                        break

                    # Stop if indentation breaks (ignoring empty lines)
                    if curr_stripped and (
                        len(curr_line) - len(curr_line.lstrip()) < start_indent
                    ):
                        break

                    bullet_lines.append(curr_line)
                    line_idx += 1

                current_blocks.append(
                    {"type": "bullet", "content": "\n".join(bullet_lines)}
                )
                continue

            # Tables
            is_table = False
            if "|" in line:
                if (line_idx + 1 < len(lines)) and re.match(
                    r"^\s*\|?.*--.*\|?\s*$", lines[line_idx + 1]
                ):
                    is_table = True

            if is_table:
                table_lines = []
                while line_idx < len(lines) and "|" in lines[line_idx]:
                    table_lines.append(lines[line_idx])
                    line_idx += 1
                current_blocks.append(
                    {"type": "table", "content": "\n".join(table_lines)}
                )
                continue

            # Plain text
            text_lines = []
            while line_idx < len(lines):
                current_line = lines[line_idx]
                curr_stripped = current_line.strip()

                # --- CRITICAL FIX 2: Stop if we hit a marker ---
                # Also removed the strict 'in list' check for flexibility
                if (
                    not curr_stripped
                    or current_line.lstrip().startswith(("-", "*", "+"))
                    or re.match(r"^\s*!\[.*\]\(.*\)\s*$", current_line)
                    or curr_stripped.startswith("```")
                    or curr_stripped.startswith(":::")  # Stop on ANY marker
                    or (
                        "|" in current_line
                        and (line_idx + 1 < len(lines))
                        and re.match(r"^\s*\|?.*--.*\|?\s*$", lines[line_idx + 1])
                    )
                ):
                    break
                text_lines.append(current_line)
                line_idx += 1

            if text_lines:
                current_blocks.append(
                    {"type": "text", "content": "\n".join(text_lines)}
                )

        slides.append(slide_data)
    return slides


# Maximum number of lines (approximated) allowed per slide.
# This is a heuristic and may need adjustment.
MAX_LINES_PER_SLIDE = 15


def get_block_len(block: Dict) -> int:
    """Estimates the 'length' of a content block in approximate lines."""
    content = block.get("content")
    if not content:
        return 0

    # Different block types contribute differently to slide 'fullness'.
    # These are heuristics and can be adjusted.
    block_type = block.get("type")
    if block_type == "text":
        # Count lines, and add a bit extra for paragraph spacing.
        return content.count("\n") + 1
    elif block_type == "bullet":
        # A rough estimate based on the number of list items.
        return len(re.findall(r"^\s*[-*+]\s", content, re.MULTILINE))
    elif block_type == "code":
        return content.count("\n") + 1
    elif block_type == "table":
        # Each table row takes up space.
        return content.count("\n") + 1
    elif block_type == "image":
        # An image is estimated to take up significant vertical space.
        return 8
    return 1


def _split_text_block(content: str, limit: int) -> List[Dict]:
    """Splits a single large text block into smaller ones by paragraph."""
    paragraphs = content.split("\n\n")
    if not any(p.strip() for p in paragraphs):
        return []

    chunks = []
    current_chunk_lines = []
    current_len = 0
    for p in paragraphs:
        p_len = p.count("\n") + 1
        if current_len + p_len > limit and current_chunk_lines:
            chunks.append("\n\n".join(current_chunk_lines))
            current_chunk_lines = [p]
            current_len = p_len
        else:
            current_chunk_lines.append(p)
            current_len += p_len

    if current_chunk_lines:
        chunks.append("\n\n".join(current_chunk_lines))

    return [{"type": "text", "content": chunk} for chunk in chunks if chunk.strip()]


def _split_bullet_block(content: str, limit: int) -> List[Dict]:
    """Splits a large bulleted list into smaller ones based on top-level items."""
    lines = content.split("\n")
    if not lines:
        return []

    min_indent = float("inf")
    for line in lines:
        if line.strip():
            min_indent = min(min_indent, len(line) - len(line.lstrip()))

    items = []
    current_item_lines = []
    for line in lines:
        if line.strip():
            indent = len(line) - len(line.lstrip())
            if indent == min_indent and line.lstrip().startswith(("-", "*", "+")):
                if current_item_lines:
                    items.append("\n".join(current_item_lines))
                current_item_lines = [line]
            else:
                current_item_lines.append(line)
        elif current_item_lines:
            current_item_lines.append(line)

    if current_item_lines:
        items.append("\n".join(current_item_lines))

    chunks = []
    current_chunk = []
    current_len = 0
    for item in items:
        item_len = item.count("\n") + 1
        if current_len + item_len > limit and current_chunk:
            chunks.append("\n".join(current_chunk))
            current_chunk = [item]
            current_len = item_len
        else:
            current_chunk.append(item)
            current_len += item_len

    if current_chunk:
        chunks.append("\n".join(current_chunk))

    return [{"type": "bullet", "content": chunk} for chunk in chunks if chunk.strip()]


def _split_code_block(content: str, language: str, limit: int) -> List[Dict]:
    """Splits a large code block into smaller ones."""
    lines = content.split("\n")
    chunks = ["\n".join(lines[i : i + limit]) for i in range(0, len(lines), limit)]
    return [
        {"type": "code", "language": language, "content": chunk}
        for chunk in chunks
        if chunk.strip()
    ]


def split_long_slides(slides_data: List[Dict]) -> List[Dict]:
    """
    Splits slides with content exceeding MAX_LINES_PER_SLIDE into multiple slides.
    This helps manage slides with a large amount of content.
    """
    new_slides_data = []
    for slide in slides_data:
        # Skip complex layouts like two-column for now to avoid overly complex logic.
        if slide.get("columns"):
            new_slides_data.append(slide)
            continue

        # 1. Expand any single block that is too long into multiple, smaller blocks.
        expanded_blocks = []
        for block in slide["blocks"]:
            # A single block (e.g., a long paragraph) might need to be split.
            if get_block_len(block) > MAX_LINES_PER_SLIDE:
                if block["type"] == "text":
                    expanded_blocks.extend(
                        _split_text_block(block["content"], MAX_LINES_PER_SLIDE)
                    )
                elif block["type"] == "bullet":
                    expanded_blocks.extend(
                        _split_bullet_block(block["content"], MAX_LINES_PER_SLIDE)
                    )
                elif block["type"] == "code":
                    expanded_blocks.extend(
                        _split_code_block(
                            block["content"],
                            block.get("language", ""),
                            MAX_LINES_PER_SLIDE,
                        )
                    )
                else:
                    # Non-splittable blocks (images, tables) are kept as-is.
                    expanded_blocks.append(block)
            else:
                expanded_blocks.append(block)

        # 2. Group the potentially expanded blocks into new slides.
        if not expanded_blocks:
            new_slides_data.append(slide)
            continue

        slide_count = 0
        current_slide_blocks = []
        current_slide_len = 0

        for block in expanded_blocks:
            block_len = get_block_len(block)
            if (
                current_slide_len + block_len > MAX_LINES_PER_SLIDE
                and current_slide_blocks
            ):
                is_first = slide_count == 0
                title = slide["title"] if is_first else f"{slide['title']} (cont.)"
                new_slides_data.append(
                    {
                        "layout": slide["layout"],
                        "title": title,
                        "blocks": current_slide_blocks,
                        "columns": [],
                        "notes": slide["notes"] if is_first else "",
                    }
                )
                slide_count += 1
                current_slide_blocks = [block]
                current_slide_len = block_len
            else:
                current_slide_blocks.append(block)
                current_slide_len += block_len

        # Add the last accumulated slide.
        if current_slide_blocks:
            is_first = slide_count == 0
            title = slide["title"] if is_first else f"{slide['title']} (cont.)"
            new_slides_data.append(
                {
                    "layout": slide["layout"],
                    "title": title,
                    "blocks": current_slide_blocks,
                    "columns": [],
                    "notes": slide["notes"] if is_first else "",
                }
            )
    return new_slides_data


def create_presentation_from_markdown(
    content: str, output_path: str = "output.pptx"
) -> str:
    """
    Creates and saves a PowerPoint presentation from structured slide data.

    This function iterates through a list of slide data dictionaries, each
    representing a slide. For each slide, it selects the appropriate layout,
    adds a new slide to the presentation, and then populates the title, notes,
    and content placeholders based on the parsed data. It contains specific
    logic to handle different layouts, such as 'two_content', and different
    content block types like text, lists, code, and tables.
    """
    prs = Presentation()
    from pptx.util import Pt

    # Map layout names from our Markdown format to the default slide layouts
    # available in python-pptx. The numbers correspond to standard layout indices.
    layout_map = {
        "title_slide": prs.slide_layouts[0],
        "title_content": prs.slide_layouts[1],
        "section_header": prs.slide_layouts[2],
        "two_content": prs.slide_layouts[
            3
        ],  # in parser we don't distinguish two_content and comparison
        "comparison": prs.slide_layouts[3],  # quick fix: use two_content for comparison
        "title_only": prs.slide_layouts[5],
        "blank": prs.slide_layouts[6],
        "picture_and_caption": prs.slide_layouts[8],
    }

    slides_data = parse_markdown_to_slides(content)
    slides_data = split_long_slides(slides_data)

    for slide_data in slides_data:
        layout_name = slide_data["layout"]
        slide_layout = layout_map.get(layout_name, layout_map["title_content"])

        current_slide = prs.slides.add_slide(slide_layout)

        # Populate the title placeholder if it exists on the layout.
        if slide_data["title"]:
            if current_slide.shapes.title:
                current_slide.shapes.title.text = slide_data["title"]

        # Populate the speaker notes.
        if slide_data.get("notes"):
            notes_slide = current_slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data["notes"]

        image_blocks = [b for b in slide_data["blocks"] if b["type"] == "image"]
        other_blocks = [b for b in slide_data["blocks"] if b["type"] != "image"]

        # --- Two-Column Layout Rendering ---
        if layout_name in ["comparison", "two_content"]:
            # Identify the left and right content placeholders.
            content_placeholders = [
                p
                for p in current_slide.placeholders
                if p.placeholder_format.idx > 0 and p.has_text_frame
            ]
            if len(content_placeholders) >= 2:
                left_ph, right_ph = content_placeholders[0], content_placeholders[1]

                # If the parser created a 'columns' data structure, render from it.
                if slide_data["columns"]:
                    # Render the first column's blocks into the left placeholder.
                    if len(slide_data["columns"]) > 0:
                        tf_left = left_ph.text_frame
                        tf_left.clear()
                        for block in slide_data["columns"][0]:
                            if block["type"] == "text":
                                p = tf_left.add_paragraph()
                                add_formatted_text_runs(p, block["content"])
                            elif block["type"] == "bullet":
                                add_bullet_points_from_markdown(
                                    tf_left, block["content"]
                                )
                            elif block["type"] == "code":
                                p = tf_left.add_paragraph()
                                run = p.add_run()
                                run.text = block["content"]
                                run.font.name = "Courier New"
                                run.font.size = Pt(10)

                    # Render the second column's blocks into the right placeholder.
                    if len(slide_data["columns"]) > 1:
                        tf_right = right_ph.text_frame
                        tf_right.clear()
                        for block in slide_data["columns"][1]:
                            if block["type"] == "text":
                                p = tf_right.add_paragraph()
                                add_formatted_text_runs(p, block["content"])
                            elif block["type"] == "bullet":
                                add_bullet_points_from_markdown(
                                    tf_right, block["content"]
                                )
                            elif block["type"] == "code":
                                p = tf_right.add_paragraph()
                                run = p.add_run()
                                run.text = block["content"]
                                run.font.name = "Courier New"
                                run.font.size = Pt(10)

                else:  # Fallback for the simple '|||' separator syntax.
                    text_block_content = ""
                    for block in other_blocks:
                        if block["type"] == "text":
                            text_block_content = block["content"]
                            break

                    left_text, right_text = (
                        text_block_content.split("|||", 1)
                        if "|||" in text_block_content
                        else (text_block_content, "")
                    )
                    left_ph.text_frame.clear()
                    p_left = left_ph.text_frame.add_paragraph()
                    add_formatted_text_runs(p_left, left_text.strip())
                    right_ph.text_frame.clear()
                    p_right = right_ph.text_frame.add_paragraph()
                    add_formatted_text_runs(p_right, right_text.strip())

        # --- Other Layouts ---
        elif layout_name == "picture_and_caption":
            if image_blocks:
                pic_placeholder = next(
                    (
                        p
                        for p in current_slide.placeholders
                        if p.placeholder_format.type == 18
                    ),
                    None,
                )
                if pic_placeholder:
                    match = re.match(r"!\[.*\]\((.*)\)", image_blocks[0]["content"])
                    if match:
                        image_path = match.group(1)
                        if os.path.exists(image_path):
                            try:
                                pic_placeholder.insert_picture(image_path)
                            except Exception as e:
                                print(f"Could not insert image {image_path}: {e}")

            body_shape = next(
                (
                    p
                    for p in current_slide.placeholders
                    if p.placeholder_format.idx > 0 and p.placeholder_format.type != 18
                ),
                None,
            )
            if body_shape and body_shape.has_text_frame:
                tf = body_shape.text_frame
                tf.clear()
                for block in other_blocks:
                    if block["type"] == "text":
                        p = tf.add_paragraph()
                        add_formatted_text_runs(p, block["content"])
                    elif block["type"] == "bullet":
                        add_bullet_points_from_markdown(tf, block["content"])
                    elif block["type"] == "code":
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = block["content"]
                        run.font.name = "Courier New"
                        run.font.size = Pt(10)

        else:  # Default handling for single-content-area layouts.
            body_shape = next(
                (
                    shape
                    for shape in current_slide.placeholders
                    if shape.placeholder_format.idx != 0 and shape.has_text_frame
                ),
                None,
            )
            if body_shape:
                tf = body_shape.text_frame
                tf.clear()

                for block in other_blocks:
                    if block["type"] == "text":
                        p = tf.add_paragraph()
                        add_formatted_text_runs(p, block["content"])
                    elif block["type"] == "bullet":
                        add_bullet_points_from_markdown(tf, block["content"])
                    elif block["type"] == "code":
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = block["content"]
                        run.font.name = "Courier New"
                        run.font.size = Pt(10)
                    elif block["type"] == "table":
                        table_data = create_table_from_markdown(block["content"])
                        if not table_data:
                            continue

                        rows, cols = len(table_data), len(table_data[0])
                        table_shape = current_slide.shapes.add_table(
                            rows,
                            cols,
                            Inches(1),
                            Inches(2.5),
                            Inches(8),
                            Inches(0.4 * (rows + 1)),
                        )
                        table = table_shape.table

                        for r_idx, row_data in enumerate(table_data):
                            for c_idx, cell_text in enumerate(row_data):
                                if c_idx < cols:
                                    cell = table.cell(r_idx, c_idx)
                                    tf = cell.text_frame
                                    tf.clear()
                                    p = tf.add_paragraph()
                                    add_formatted_text_runs(p, cell_text)

            # Add any images to the slide at a default position.
            if image_blocks:
                for i, block in enumerate(image_blocks):
                    match = re.match(r"!\[.*\]\((.*)\)", block["content"])
                    if match:
                        image_path = match.group(1)
                        if os.path.exists(image_path):
                            left, top, height = (
                                Inches(1),
                                Inches(2.5 + i * 2),
                                Inches(2),
                            )
                            try:
                                current_slide.shapes.add_picture(
                                    image_path, left, top, height=height
                                )
                            except Exception as e:
                                print(f"Could not add image {image_path}: {e}")

    prs.save(output_path)
    return output_path


def generate_gemini_response(text: str, model: str = "gemini-2.5-pro"):
    client = genai.Client(api_key="AIzaSyA4YsTnbNjl2gKn20EqPa-9nom9yymEwd0")
    try:
        response = client.models.generate_content(
            model=model,
            contents=GEMINI_LECTURE_PROMPT + text,
        )
        logging.info(f"{response.text}")
    except Exception as e:
        print("=== oh poor thing, an error ===")
        raise e

    return response.text


# to delete
def extract_content_with_docling(
    file_path: str, enabled_ocr=True, page_range: str = None
):
    """
    Extracts content from various file types using the docling library.

    Args:
        file_path (str): The path to the file to extract content from.
        page_range (str): Optional page range (e.g., "1-5", "2,4").

    Returns:
        Dict: A dictionary containing the extracted content.
    """
    suffix = file_path.split(".")[-1]
    try:
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import (
            PdfPipelineOptions,
            TesseractCliOcrOptions,
            TesseractOcrOptions,
        )
        from docling.document_converter import DocumentConverter, PdfFormatOption

        # Set lang=["auto"] with a tesseract OCR engine: TesseractOcrOptions, TesseractCliOcrOptions
        # ocr_options = TesseractOcrOptions(lang=["auto"])

        ocr_options = TesseractCliOcrOptions(lang=["eng"])

        pipeline_options = PdfPipelineOptions(
            do_ocr=enabled_ocr, do_table_structure=True, ocr_options=ocr_options
        )

        doc_converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(
                    pipeline_options=pipeline_options,
                )
            }
        )

        doc = doc_converter.convert(file_path).document

        # Apply page range if provided
        if page_range:
            pages = []
            for part in page_range.split(","):
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    pages.extend(range(start - 1, end))  # 0-indexed
                else:
                    pages.append(int(part) - 1)  # 0-indexed

            # Filter pages
            filtered_content = []
            for i, page in enumerate(doc.pages):
                if i in pages:
                    filtered_content.append(page.export_to_markdown())
            return "\n".join(filtered_content)
        else:
            pprint(doc.export_to_markdown())
            return doc.export_to_markdown()

    except Exception as e:
        print(f"Error extracting content with docling from {file_path}: {e}")
        return {"error": str(e)}


def sleep_min(seconds=60):
    import time

    logging.info(f"sleeping for {seconds}")
    time.sleep(seconds)


if __name__ == "__main__":
    example_content = """
# My Presentation
- Main Point 1
  - Subpoint 1.1
- Main Point 2

::: notes
This is a speaker note for the first slide.
:::

---
layout: two_content
# Code and Bullets

::: column
```python
def fib(n):
    a, b = 0, 1
    while a < n:
        print(a, end=' ')
        a, b = b, a+b
    print()
```
:::

::: column
- This is a list of features.
- It is in the right column.
- Next to a code block.
:::

---
layout: picture_and_caption
# Image Slide
![A generated image](gemini:A cute robot waving)
This is the caption for the image.

---

layout: comparison
# Comparison Example
Left side content.
|||
Right side content.

---

# Table Slide
| Header 1 | Header 2 |
|----------|----------|
| Cell 1   | Cell 2   |
| Cell 3   | Cell 4   |

---

# Slide with an Image
This slide has text and an image.
![Another image](/path/to/your/other_image.png)
The image will be placed at a default position.
"""

    source_data_doc = extract_content_with_docling(
        "L2 Development of Aortic arches.pdf"
    )
    logging.info("-------------- started gemini")
    real_content = (
        generate_gemini_response(source_data_doc)
        or "# sorry error yr lovely llm model didnt generate"
    )
    print("-------------- preprocessing images")

    real_content_with_images = preprocess_markdown_for_images(real_content)
    output_path = create_presentation_from_markdown(
        real_content_with_images, "my_markdown_presentation.pptx"
    )
    print(f"Created presentation: {output_path}")

    # Create a dummy file for docling to extract from
    dummy_text_file = "dummy_docling_test.txt"
    with open(dummy_text_file, "w") as f:
        f.write("This is a test document for docling extraction.\n")

    # # Clean up the dummy file
    # import os

    # os.remove(dummy_text_file)


def generate_image(prompt: str, output_path: str):
    """Generates an image using Gemini and saves it to a file."""
    try:
        logging.info(f"Generating image for prompt: {prompt}")
        client = genai.Client(api_key="AIzaSyA4YsTnbNjl2gKn20EqPa-9nom9yymEwd0")
        response = client.models.generate_content(
            model="gemini-2.0-flash-preview-image-generation",  # Or another model, as per docs
            contents=[prompt],
        )

        for part in response.parts:
            if part.inline_data is not None:
                image = part.as_image()
                image.save(output_path)
                logging.info(f"Saved generated image to {output_path}")
                return output_path
    except Exception as e:
        logging.error(f"Error generating image: {e}")
        return None


def preprocess_markdown_for_images(markdown_content: str) -> str:
    """
    Finds image generation directives, generates images, and replaces
    the directive with the path to the generated image file.api_key="AIzaSyA4YsTnbNjl2gKn20EqPa-9nom9yymEwd0"
    """
    # Regex to find ![alt text](gemini:prompt)
    pattern = r"!\[(.*?)\]\(gemini:(.*?)\)"
    matches = re.findall(pattern, markdown_content)

    if not matches:
        return markdown_content

    images_dir = "generated_images"
    os.makedirs(images_dir, exist_ok=True)

    modified_content = markdown_content
    for i, (alt_text, prompt) in enumerate(matches):
        image_filename = f"generated_image_{i}.png"
        output_path = os.path.join(images_dir, image_filename)
        sleep_min()
        if generate_image(prompt, output_path):
            # Replace the gemini directive with the local file path
            original_directive = f"![{alt_text}](gemini:{prompt})"
            new_directive = f"![{alt_text}]({output_path})"
            modified_content = modified_content.replace(
                original_directive, new_directive, 1
            )
            logging.info(f"{new_directive = }\n {original_directive = }\n")
        else:
            logging.warning(f"Could not generate image for prompt: {prompt}")

    return modified_content


def create_table_from_markdown(text: str) -> List[List[str]]:
    """Convert Markdown table to table data."""
    lines = [line.strip() for line in text.split("\n") if line.strip()]
    if not lines:
        return []

    # Filter out separator line
    lines = [l for l in lines if not re.match(r"^\s*\|?.*--.*\|?\s*$", l)]

    table_data = []
    for row_str in lines:
        if row_str.startswith("|"):
            row_str = row_str[1:]
        if row_str.endswith("|"):
            row_str = row_str[:-1]

        cells = [cell.strip() for cell in row_str.split("|")]
        table_data.append(cells)

    return table_data


def add_formatted_text_runs(paragraph, text, bold=False, italic=False, underline=False):
    """
    Recursively parses markdown-like formatting for bold, italic, and
    underline, and adds formatted runs to the given paragraph.

    This function works by repeatedly splitting the text by all possible
    formatting markers. If the text can't be split further, it's considered
    plain text and added as a single run with the formatting inherited from
    its parent calls. If the text is split, the function calls itself for
    each segment, updating the formatting flags (bold, italic, etc.) for
    segments that were wrapped in markers.

    - ***text*** for bold & italic
    - **text** for bold
    - *text* for italic
    - __text__ for underline
    """
    # Split by the formatting markers, keeping them.
    parts = re.split(
        r"(\*\*\*[\s\S]+?\*\*\*|\*\*[\s\S]+?\*\*|\*[\s\S]+?\*|__[\s\S]+?__)", text
    )

    # Base case: If the text is not split into more than one part,
    # it means no formatting markers were found. Add the text as a single run.
    if len(parts) == 1:
        if text:
            run = paragraph.add_run()
            run.text = text
            font = run.font
            font.bold = bold
            font.italic = italic
            font.underline = underline
        return

    # Recursive step: For each part of the split text, determine if it's
    # a formatted segment or plain text, and recurse.
    for part in parts:
        if not part:
            continue

        # Order of checks is important: *** must be checked before ** or *.
        if part.startswith("***") and part.endswith("***"):
            # It's a bold and italic segment, recurse on its content with new formatting.
            add_formatted_text_runs(
                paragraph, part[3:-3], bold=True, italic=True, underline=underline
            )
        elif part.startswith("**") and part.endswith("**"):
            # It's a bold segment, recurse on its content.
            add_formatted_text_runs(
                paragraph, part[2:-2], bold=True, italic=italic, underline=underline
            )
        elif part.startswith("*") and part.endswith("*"):
            # It's an italic segment, recurse on its content.
            add_formatted_text_runs(
                paragraph, part[1:-1], bold=bold, italic=True, underline=underline
            )
        elif part.startswith("__") and part.endswith("__"):
            # It's an underline segment, recurse on its content.
            add_formatted_text_runs(
                paragraph, part[2:-2], bold=bold, italic=italic, underline=True
            )
        else:
            # This is a plain text part with no new markers.
            # Recurse on it, passing down the formatting from the parent.
            add_formatted_text_runs(
                paragraph, part, bold=bold, italic=italic, underline=underline
            )


def add_bullet_points_from_markdown(text_frame, points: str):
    """Add bullet points to a text frame from Markdown list."""
    if not text_frame.text.strip():
        text_frame.text = ""

    def get_level_and_text(line: str) -> tuple[int, str]:
        """Determine level and clean text from a Markdown list item."""
        stripped_line = line.lstrip()
        text = stripped_line

        if text.startswith(("-", "*", "+")) and text[1:2] in (" ", ""):
            text = text[1:].lstrip()

        indent = len(line) - len(line.lstrip())
        level = indent // 2
        return level, text

    lines = [line for line in points.split("\n") if line.strip()]
    if not lines:
        return

    for line in lines:
        level, text = get_level_and_text(line)
        p = text_frame.add_paragraph()
        add_formatted_text_runs(p, text)
        p.level = min(level, 8)


def parse_markdown_to_slides(content: str) -> List[Dict]:
    """
    Parses markdown-like text into a list of slide definitions.

    This function iterates through the input content line by line, maintaining
    state (e.g., `in_notes_block`, `in_column_block`) to determine what kind of
    content is currently being parsed. It identifies structural elements like
    layout definitions, titles, and various content blocks (notes, columns,
    code, lists, tables, text), converting them into a structured list of
    dictionaries, where each dictionary represents a single slide.
    """
    slides = []
    # A slide is defined as the content between '---' separators.
    import re


from typing import Dict, List


def parse_markdown_to_slides(content: str) -> List[Dict]:
    slides = []
    slide_contents = re.split(r"\n---\n", content)

    for slide_content in slide_contents:
        if not slide_content.strip():
            continue

        lines = slide_content.strip().split("\n")

        slide_data = {
            "layout": "title_content",
            "title": None,
            "blocks": [],
            "columns": [],
            "notes": "",
        }

        # The first lines of a slide can define the layout and title.
        if lines and lines[0].startswith("layout:"):
            slide_data["layout"] = lines[0].split(":", 1)[1].strip()
            lines.pop(0)

        if lines and lines[0].startswith("# "):
            slide_data["title"] = lines[0][2:].strip()
            lines.pop(0)

        line_idx = 0
        in_notes_block = False
        in_column_block = False

        while line_idx < len(lines):
            line = lines[line_idx]
            stripped_line = line.strip()  # Calculate this once

            # 1. ROBUST MARKER CHECKS (Handle extra spaces/case)
            if stripped_line.lower().startswith("::: notes"):
                in_notes_block = True
                line_idx += 1
                continue

            if stripped_line.lower().startswith("::: column"):
                in_column_block = True
                if not slide_data["columns"] and slide_data["blocks"]:
                    slide_data["columns"].append(slide_data["blocks"])
                    slide_data["blocks"] = []
                slide_data["columns"].append([])
                line_idx += 1
                continue

            if stripped_line == ":::":
                in_notes_block = False
                in_column_block = False
                line_idx += 1
                continue

            if in_notes_block:
                slide_data["notes"] += line + "\n"
                line_idx += 1
                continue

            current_blocks = (
                slide_data["columns"][-1]
                if in_column_block and slide_data["columns"]
                else slide_data["blocks"]
            )

            if not stripped_line:
                line_idx += 1
                continue

            # --- PARSING BLOCKS ---

            # Code blocks
            if stripped_line.startswith("```"):
                code_lines = []
                lang = stripped_line[3:]
                line_idx += 1
                while line_idx < len(lines) and not lines[line_idx].strip().startswith(
                    "```"
                ):
                    code_lines.append(lines[line_idx])
                    line_idx += 1
                line_idx += 1
                current_blocks.append(
                    {"type": "code", "language": lang, "content": "\n".join(code_lines)}
                )
                continue

            # Image
            if re.match(r"^\s*!\[.*\]\(.*\)\s*$", line):
                current_blocks.append({"type": "image", "content": line.strip()})
                line_idx += 1
                continue

            # Bulleted lists
            if line.lstrip().startswith(("-", "*", "+")):
                bullet_lines = []
                start_indent = len(line) - len(line.lstrip())

                while line_idx < len(lines):
                    curr_line = lines[line_idx]
                    curr_stripped = curr_line.strip()

                    # --- CRITICAL FIX 1: Stop if we hit a marker ---
                    if curr_stripped.startswith(":::") or curr_stripped.startswith(
                        "```"
                    ):
                        break

                    # Stop if indentation breaks (ignoring empty lines)
                    if curr_stripped and (
                        len(curr_line) - len(curr_line.lstrip()) < start_indent
                    ):
                        break

                    bullet_lines.append(curr_line)
                    line_idx += 1

                current_blocks.append(
                    {"type": "bullet", "content": "\n".join(bullet_lines)}
                )
                continue

            # Tables
            is_table = False
            if "|" in line:
                if (line_idx + 1 < len(lines)) and re.match(
                    r"^\s*\|?.*--.*\|?\s*$", lines[line_idx + 1]
                ):
                    is_table = True

            if is_table:
                table_lines = []
                while line_idx < len(lines) and "|" in lines[line_idx]:
                    table_lines.append(lines[line_idx])
                    line_idx += 1
                current_blocks.append(
                    {"type": "table", "content": "\n".join(table_lines)}
                )
                continue

            # Plain text
            text_lines = []
            while line_idx < len(lines):
                current_line = lines[line_idx]
                curr_stripped = current_line.strip()

                # --- CRITICAL FIX 2: Stop if we hit a marker ---
                # Also removed the strict 'in list' check for flexibility
                if (
                    not curr_stripped
                    or current_line.lstrip().startswith(("-", "*", "+"))
                    or re.match(r"^\s*!\[.*\]\(.*\)\s*$", current_line)
                    or curr_stripped.startswith("```")
                    or curr_stripped.startswith(":::")  # Stop on ANY marker
                    or (
                        "|" in current_line
                        and (line_idx + 1 < len(lines))
                        and re.match(r"^\s*\|?.*--.*\|?\s*$", lines[line_idx + 1])
                    )
                ):
                    break
                text_lines.append(current_line)
                line_idx += 1

            if text_lines:
                current_blocks.append(
                    {"type": "text", "content": "\n".join(text_lines)}
                )

        slides.append(slide_data)
    return slides


# Maximum number of lines (approximated) allowed per slide.
# This is a heuristic and may need adjustment.
MAX_LINES_PER_SLIDE = 15


def get_block_len(block: Dict) -> int:
    """Estimates the 'length' of a content block in approximate lines."""
    content = block.get("content")
    if not content:
        return 0

    # Different block types contribute differently to slide 'fullness'.
    # These are heuristics and can be adjusted.
    block_type = block.get("type")
    if block_type == "text":
        # Count lines, and add a bit extra for paragraph spacing.
        return content.count("\n") + 1
    elif block_type == "bullet":
        # A rough estimate based on the number of list items.
        return len(re.findall(r"^\s*[-*+]\s", content, re.MULTILINE))
    elif block_type == "code":
        return content.count("\n") + 1
    elif block_type == "table":
        # Each table row takes up space.
        return content.count("\n") + 1
    elif block_type == "image":
        # An image is estimated to take up significant vertical space.
        return 8
    return 1


def _split_text_block(content: str, limit: int) -> List[Dict]:
    """Splits a single large text block into smaller ones by paragraph."""
    paragraphs = content.split("\n\n")
    if not any(p.strip() for p in paragraphs):
        return []

    chunks = []
    current_chunk_lines = []
    current_len = 0
    for p in paragraphs:
        p_len = p.count("\n") + 1
        if current_len + p_len > limit and current_chunk_lines:
            chunks.append("\n\n".join(current_chunk_lines))
            current_chunk_lines = [p]
            current_len = p_len
        else:
            current_chunk_lines.append(p)
            current_len += p_len

    if current_chunk_lines:
        chunks.append("\n\n".join(current_chunk_lines))

    return [{"type": "text", "content": chunk} for chunk in chunks if chunk.strip()]


def _split_bullet_block(content: str, limit: int) -> List[Dict]:
    """Splits a large bulleted list into smaller ones based on top-level items."""
    lines = content.split("\n")
    if not lines:
        return []

    min_indent = float("inf")
    for line in lines:
        if line.strip():
            min_indent = min(min_indent, len(line) - len(line.lstrip()))

    items = []
    current_item_lines = []
    for line in lines:
        if line.strip():
            indent = len(line) - len(line.lstrip())
            if indent == min_indent and line.lstrip().startswith(("-", "*", "+")):
                if current_item_lines:
                    items.append("\n".join(current_item_lines))
                current_item_lines = [line]
            else:
                current_item_lines.append(line)
        elif current_item_lines:
            current_item_lines.append(line)

    if current_item_lines:
        items.append("\n".join(current_item_lines))

    chunks = []
    current_chunk = []
    current_len = 0
    for item in items:
        item_len = item.count("\n") + 1
        if current_len + item_len > limit and current_chunk:
            chunks.append("\n".join(current_chunk))
            current_chunk = [item]
            current_len = item_len
        else:
            current_chunk.append(item)
            current_len += item_len

    if current_chunk:
        chunks.append("\n".join(current_chunk))

    return [{"type": "bullet", "content": chunk} for chunk in chunks if chunk.strip()]


def _split_code_block(content: str, language: str, limit: int) -> List[Dict]:
    """Splits a large code block into smaller ones."""
    lines = content.split("\n")
    chunks = ["\n".join(lines[i : i + limit]) for i in range(0, len(lines), limit)]
    return [
        {"type": "code", "language": language, "content": chunk}
        for chunk in chunks
        if chunk.strip()
    ]


def split_long_slides(slides_data: List[Dict]) -> List[Dict]:
    """
    Splits slides with content exceeding MAX_LINES_PER_SLIDE into multiple slides.
    This helps manage slides with a large amount of content.
    """
    new_slides_data = []
    for slide in slides_data:
        # Skip complex layouts like two-column for now to avoid overly complex logic.
        if slide.get("columns"):
            new_slides_data.append(slide)
            continue

        # 1. Expand any single block that is too long into multiple, smaller blocks.
        expanded_blocks = []
        for block in slide["blocks"]:
            # A single block (e.g., a long paragraph) might need to be split.
            if get_block_len(block) > MAX_LINES_PER_SLIDE:
                if block["type"] == "text":
                    expanded_blocks.extend(
                        _split_text_block(block["content"], MAX_LINES_PER_SLIDE)
                    )
                elif block["type"] == "bullet":
                    expanded_blocks.extend(
                        _split_bullet_block(block["content"], MAX_LINES_PER_SLIDE)
                    )
                elif block["type"] == "code":
                    expanded_blocks.extend(
                        _split_code_block(
                            block["content"],
                            block.get("language", ""),
                            MAX_LINES_PER_SLIDE,
                        )
                    )
                else:
                    # Non-splittable blocks (images, tables) are kept as-is.
                    expanded_blocks.append(block)
            else:
                expanded_blocks.append(block)

        # 2. Group the potentially expanded blocks into new slides.
        if not expanded_blocks:
            new_slides_data.append(slide)
            continue

        slide_count = 0
        current_slide_blocks = []
        current_slide_len = 0

        for block in expanded_blocks:
            block_len = get_block_len(block)
            if (
                current_slide_len + block_len > MAX_LINES_PER_SLIDE
                and current_slide_blocks
            ):
                is_first = slide_count == 0
                title = slide["title"] if is_first else f"{slide['title']} (cont.)"
                new_slides_data.append(
                    {
                        "layout": slide["layout"],
                        "title": title,
                        "blocks": current_slide_blocks,
                        "columns": [],
                        "notes": slide["notes"] if is_first else "",
                    }
                )
                slide_count += 1
                current_slide_blocks = [block]
                current_slide_len = block_len
            else:
                current_slide_blocks.append(block)
                current_slide_len += block_len

        # Add the last accumulated slide.
        if current_slide_blocks:
            is_first = slide_count == 0
            title = slide["title"] if is_first else f"{slide['title']} (cont.)"
            new_slides_data.append(
                {
                    "layout": slide["layout"],
                    "title": title,
                    "blocks": current_slide_blocks,
                    "columns": [],
                    "notes": slide["notes"] if is_first else "",
                }
            )
    return new_slides_data


def create_presentation_from_markdown(
    content: str, output_path: str = "output.pptx"
) -> str:
    """
    Creates and saves a PowerPoint presentation from structured slide data.

    This function iterates through a list of slide data dictionaries, each
    representing a slide. For each slide, it selects the appropriate layout,
    adds a new slide to the presentation, and then populates the title, notes,
    and content placeholders based on the parsed data. It contains specific
    logic to handle different layouts, such as 'two_content', and different
    content block types like text, lists, code, and tables.
    """
    prs = Presentation()
    from pptx.util import Pt

    # Map layout names from our Markdown format to the default slide layouts
    # available in python-pptx. The numbers correspond to standard layout indices.
    layout_map = {
        "title_slide": prs.slide_layouts[0],
        "title_content": prs.slide_layouts[1],
        "section_header": prs.slide_layouts[2],
        "two_content": prs.slide_layouts[
            3
        ],  # in parser we don't distinguish two_content and comparison
        "comparison": prs.slide_layouts[3],  # quick fix: use two_content for comparison
        "title_only": prs.slide_layouts[5],
        "blank": prs.slide_layouts[6],
        "picture_and_caption": prs.slide_layouts[8],
    }

    slides_data = parse_markdown_to_slides(content)
    slides_data = split_long_slides(slides_data)

    for slide_data in slides_data:
        layout_name = slide_data["layout"]
        slide_layout = layout_map.get(layout_name, layout_map["title_content"])

        current_slide = prs.slides.add_slide(slide_layout)

        # Populate the title placeholder if it exists on the layout.
        if slide_data["title"]:
            if current_slide.shapes.title:
                current_slide.shapes.title.text = slide_data["title"]

        # Populate the speaker notes.
        if slide_data.get("notes"):
            notes_slide = current_slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data["notes"]

        image_blocks = [b for b in slide_data["blocks"] if b["type"] == "image"]
        other_blocks = [b for b in slide_data["blocks"] if b["type"] != "image"]

        # --- Two-Column Layout Rendering ---
        if layout_name in ["comparison", "two_content"]:
            # Identify the left and right content placeholders.
            content_placeholders = [
                p
                for p in current_slide.placeholders
                if p.placeholder_format.idx > 0 and p.has_text_frame
            ]
            if len(content_placeholders) >= 2:
                left_ph, right_ph = content_placeholders[0], content_placeholders[1]

                # If the parser created a 'columns' data structure, render from it.
                if slide_data["columns"]:
                    # Render the first column's blocks into the left placeholder.
                    if len(slide_data["columns"]) > 0:
                        tf_left = left_ph.text_frame
                        tf_left.clear()
                        for block in slide_data["columns"][0]:
                            if block["type"] == "text":
                                p = tf_left.add_paragraph()
                                add_formatted_text_runs(p, block["content"])
                            elif block["type"] == "bullet":
                                add_bullet_points_from_markdown(
                                    tf_left, block["content"]
                                )
                            elif block["type"] == "code":
                                p = tf_left.add_paragraph()
                                run = p.add_run()
                                run.text = block["content"]
                                run.font.name = "Courier New"
                                run.font.size = Pt(10)

                    # Render the second column's blocks into the right placeholder.
                    if len(slide_data["columns"]) > 1:
                        tf_right = right_ph.text_frame
                        tf_right.clear()
                        for block in slide_data["columns"][1]:
                            if block["type"] == "text":
                                p = tf_right.add_paragraph()
                                add_formatted_text_runs(p, block["content"])
                            elif block["type"] == "bullet":
                                add_bullet_points_from_markdown(
                                    tf_right, block["content"]
                                )
                            elif block["type"] == "code":
                                p = tf_right.add_paragraph()
                                run = p.add_run()
                                run.text = block["content"]
                                run.font.name = "Courier New"
                                run.font.size = Pt(10)

                else:  # Fallback for the simple '|||' separator syntax.
                    text_block_content = ""
                    for block in other_blocks:
                        if block["type"] == "text":
                            text_block_content = block["content"]
                            break

                    left_text, right_text = (
                        text_block_content.split("|||", 1)
                        if "|||" in text_block_content
                        else (text_block_content, "")
                    )
                    left_ph.text_frame.clear()
                    p_left = left_ph.text_frame.add_paragraph()
                    add_formatted_text_runs(p_left, left_text.strip())
                    right_ph.text_frame.clear()
                    p_right = right_ph.text_frame.add_paragraph()
                    add_formatted_text_runs(p_right, right_text.strip())

        # --- Other Layouts ---
        elif layout_name == "picture_and_caption":
            if image_blocks:
                pic_placeholder = next(
                    (
                        p
                        for p in current_slide.placeholders
                        if p.placeholder_format.type == 18
                    ),
                    None,
                )
                if pic_placeholder:
                    match = re.match(r"!\[.*\]\((.*)\)", image_blocks[0]["content"])
                    if match:
                        image_path = match.group(1)
                        if os.path.exists(image_path):
                            try:
                                pic_placeholder.insert_picture(image_path)
                            except Exception as e:
                                print(f"Could not insert image {image_path}: {e}")

            body_shape = next(
                (
                    p
                    for p in current_slide.placeholders
                    if p.placeholder_format.idx > 0 and p.placeholder_format.type != 18
                ),
                None,
            )
            if body_shape and body_shape.has_text_frame:
                tf = body_shape.text_frame
                tf.clear()
                for block in other_blocks:
                    if block["type"] == "text":
                        p = tf.add_paragraph()
                        add_formatted_text_runs(p, block["content"])
                    elif block["type"] == "bullet":
                        add_bullet_points_from_markdown(tf, block["content"])
                    elif block["type"] == "code":
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = block["content"]
                        run.font.name = "Courier New"
                        run.font.size = Pt(10)

        else:  # Default handling for single-content-area layouts.
            body_shape = next(
                (
                    shape
                    for shape in current_slide.placeholders
                    if shape.placeholder_format.idx != 0 and shape.has_text_frame
                ),
                None,
            )
            if body_shape:
                tf = body_shape.text_frame
                tf.clear()

                for block in other_blocks:
                    if block["type"] == "text":
                        p = tf.add_paragraph()
                        add_formatted_text_runs(p, block["content"])
                    elif block["type"] == "bullet":
                        add_bullet_points_from_markdown(tf, block["content"])
                    elif block["type"] == "code":
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = block["content"]
                        run.font.name = "Courier New"
                        run.font.size = Pt(10)
                    elif block["type"] == "table":
                        table_data = create_table_from_markdown(block["content"])
                        if not table_data:
                            continue

                        rows, cols = len(table_data), len(table_data[0])
                        table_shape = current_slide.shapes.add_table(
                            rows,
                            cols,
                            Inches(1),
                            Inches(2.5),
                            Inches(8),
                            Inches(0.4 * (rows + 1)),
                        )
                        table = table_shape.table

                        for r_idx, row_data in enumerate(table_data):
                            for c_idx, cell_text in enumerate(row_data):
                                if c_idx < cols:
                                    cell = table.cell(r_idx, c_idx)
                                    tf = cell.text_frame
                                    tf.clear()
                                    p = tf.add_paragraph()
                                    add_formatted_text_runs(p, cell_text)

            # Add any images to the slide at a default position.
            if image_blocks:
                for i, block in enumerate(image_blocks):
                    match = re.match(r"!\[.*\]\((.*)\)", block["content"])
                    if match:
                        image_path = match.group(1)
                        if os.path.exists(image_path):
                            left, top, height = (
                                Inches(1),
                                Inches(2.5 + i * 2),
                                Inches(2),
                            )
                            try:
                                current_slide.shapes.add_picture(
                                    image_path, left, top, height=height
                                )
                            except Exception as e:
                                print(f"Could not add image {image_path}: {e}")

    prs.save(output_path)
    return output_path


def generate_gemini_response(text: str, model: str = "gemini-2.5-pro"):
    client = genai.Client(api_key="AIzaSyA4YsTnbNjl2gKn20EqPa-9nom9yymEwd0")
    try:
        response = client.models.generate_content(
            model=model,
            contents=GEMINI_LECTURE_PROMPT + text,
        )
        logging.info(f"{response.text}")
    except Exception as e:
        print("=== oh poor thing, an error ===")
        raise e

    return response.text


# to delete
def extract_content_with_docling(
    file_path: str, enabled_ocr=True, page_range: str = None
):
    """
    Extracts content from various file types using the docling library.

    Args:
        file_path (str): The path to the file to extract content from.
        page_range (str): Optional page range (e.g., "1-5", "2,4").

    Returns:
        Dict: A dictionary containing the extracted content.
    """
    suffix = file_path.split(".")[-1]
    try:
        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import (
            PdfPipelineOptions,
            TesseractCliOcrOptions,
            TesseractOcrOptions,
        )
        from docling.document_converter import DocumentConverter, PdfFormatOption

        # Set lang=["auto"] with a tesseract OCR engine: TesseractOcrOptions, TesseractCliOcrOptions
        # ocr_options = TesseractOcrOptions(lang=["auto"])

        ocr_options = TesseractCliOcrOptions(lang=["eng"])

        pipeline_options = PdfPipelineOptions(
            do_ocr=enabled_ocr, do_table_structure=True, ocr_options=ocr_options
        )

        doc_converter = DocumentConverter(
            format_options={
                InputFormat.PDF: PdfFormatOption(
                    pipeline_options=pipeline_options,
                )
            }
        )

        doc = doc_converter.convert(file_path).document

        # Apply page range if provided
        if page_range:
            pages = []
            for part in page_range.split(","):
                if "-" in part:
                    start, end = map(int, part.split("-"))
                    pages.extend(range(start - 1, end))  # 0-indexed
                else:
                    pages.append(int(part) - 1)  # 0-indexed

            # Filter pages
            filtered_content = []
            for i, page in enumerate(doc.pages):
                if i in pages:
                    filtered_content.append(page.export_to_markdown())
            return "\n".join(filtered_content)
        else:
            pprint(doc.export_to_markdown())
            return doc.export_to_markdown()

    except Exception as e:
        print(f"Error extracting content with docling from {file_path}: {e}")
        return {"error": str(e)}


def sleep_min(seconds=60):
    import time

    logging.info(f"sleeping for {seconds}")
    time.sleep(seconds)


if __name__ == "__main__":
    example_content = """
# My Presentation
- Main Point 1
  - Subpoint 1.1
- Main Point 2

::: notes
This is a speaker note for the first slide.
:::

---
layout: two_content
# Code and Bullets

::: column
```python
def fib(n):
    a, b = 0, 1
    while a < n:
        print(a, end=' ')
        a, b = b, a+b
    print()
```
:::

::: column
- This is a list of features.
- It is in the right column.
- Next to a code block.
:::

---
layout: picture_and_caption
# Image Slide
![A generated image](gemini:A cute robot waving)
This is the caption for the image.

---

layout: comparison
# Comparison Example
Left side content.
|||
Right side content.

---

# Table Slide
| Header 1 | Header 2 |
|----------|----------|
| Cell 1   | Cell 2   |
| Cell 3   | Cell 4   |

---

# Slide with an Image
This slide has text and an image.
![Another image](/path/to/your/other_image.png)
The image will be placed at a default position.
"""

    source_data_doc = extract_content_with_docling(
        "L2 Development of Aortic arches.pdf"
    )
    logging.info("-------------- started gemini")
    real_content = (
        generate_gemini_response(source_data_doc)
        or "# sorry error yr lovely llm model didnt generate"
    )
    print("-------------- preprocessing images")

    real_content_with_images = preprocess_markdown_for_images(real_content)
    output_path = create_presentation_from_markdown(
        real_content_with_images, "my_markdown_presentation.pptx"
    )
    print(f"Created presentation: {output_path}")

    # Create a dummy file for docling to extract from
    dummy_text_file = "dummy_docling_test.txt"
    with open(dummy_text_file, "w") as f:
        f.write("This is a test document for docling extraction.\n")

    # # Clean up the dummy file
    # import os

    # os.remove(dummy_text_file)
