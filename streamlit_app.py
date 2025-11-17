import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from typing import List, Dict
import re
import os
import tempfile
import httpx

# --- Dependency Check and Mocking ---
# IMPORTANT: For full functionality, the 'docling' library and its dependencies
# (including Tesseract OCR) must be installed in the Streamlit environment.
try:
    from docling.document_converter import DocumentConverter
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import (
        PdfPipelineOptions,
        TesseractCliOcrOptions,
        TesseractOcrOptions,
    )
    from docling.document_converter import PdfFormatOption
    DOCLING_AVAILABLE = True
except ImportError:
    st.warning("`docling` library not found. Extraction features will be mocked. Please install `docling` and its dependencies (like Tesseract) for full functionality.")
    DOCLING_AVAILABLE = False
    # Mock function definition for unavailable docling
    def extract_content_with_docling(file_path: str, page_range: str = None) :
        return f"# Extracted Content from {os.path.basename(file_path)}\n\nThis is mock content because docling is not installed or available.\n\n- Main Topic 1\n  - Sub Point A\n- Main Topic 2\n\n---\n\n# Slide Title 2\n| Col 1 | Col 2 |\n|---|---|\n| Data 1 | Data 2 |"

# --- Constants and Utility Functions from pptx_utils.py ---

GEMINI_LECTURE_PROMPT = '''
You are an expert instructional designer and content synthesizer. Your task is to take the provided lecture text and convert it into a structured JSON object representing a slide presentation...
'''
# NOTE: The genai functions are omitted/commented out as they require API key setup, which is outside this conversion task's scope.

def create_table_from_markdown(text: str) -> List[List[str]]:
    """Convert Markdown table to table data."""
    lines = [line.strip() for line in text.split('\n') if line.strip()]
    if not lines: return []
    lines = [l for l in lines if not re.match(r'^\s*\|?.*--.*\|?\s*$', l)]
    
    table_data = []
    for row_str in lines:
        row_str = row_str.strip()
        if row_str.startswith('|'): row_str = row_str[1:]
        if row_str.endswith('|'): row_str = row_str[:-1]
        cells = [cell.strip() for cell in row_str.split('|')]
        table_data.append(cells)
    return table_data

def add_bullet_points_from_markdown(text_frame, points: str):
    """Add bullet points to a text frame from Markdown list."""
    if not text_frame.text.strip(): text_frame.text = ""

    def get_level_and_text(line: str) -> tuple[int, str]:
        stripped_line = line.lstrip()
        text = stripped_line
        if text.startswith(('-', '*', '+')) and text[1:2] in (' ', ''):
            text = text[1:].lstrip()
        indent = len(line) - len(line.lstrip())
        level = indent // 2
        return level, text

    lines = [line for line in points.split('\n') if line.strip()]
    if not lines: return

    for line in lines:
        level, text = get_level_and_text(line)
        p = text_frame.add_paragraph()
        p.text = text
        p.level = min(level, 8)

def parse_markdown_to_slides(content: str) -> List[Dict]:
    """Parses markdown-like text into a list of slide definitions."""
    slides = []
    slide_contents = re.split(r'\n---\n', content)

    for slide_content in slide_contents:
        if not slide_content.strip(): continue

        lines = slide_content.strip().split('\n')
        slide_data = {'layout': 'title_content', 'title': None, 'blocks': []}

        # Check for layout override
        if lines and lines[0].strip().lower().startswith('layout:'):
            slide_data['layout'] = lines[0].split(':', 1)[1].strip()
            lines.pop(0)
        
        # Check for title (H1)
        if lines and lines[0].strip().startswith('# '):
            slide_data['title'] = lines[0][2:].strip()
            lines.pop(0)

        line_idx = 0
        while line_idx < len(lines):
            line = lines[line_idx]
            if not line.strip(): line_idx += 1; continue

            # Bullet points block
            if line.lstrip().startswith(('-', '*', '+')):
                bullet_lines = []
                start_indent = len(line) - len(line.lstrip())
                while line_idx < len(lines):
                    current_line = lines[line_idx]
                    current_indent = len(current_line) - len(current_line.lstrip())
                    
                    if current_line.strip() and current_indent < start_indent and not current_line.lstrip().startswith(('-', '*', '+')):
                        break # Unindented line is not part of the list
                    
                    if not current_line.strip() and line_idx + 1 < len(lines) and lines[line_idx+1].strip() and (len(lines[line_idx+1]) - len(lines[line_idx+1].lstrip()) < start_indent):
                        break # Empty line followed by unindented line
                        
                    if current_line.strip() or current_indent >= start_indent:
                        bullet_lines.append(current_line)
                    line_idx += 1
                
                slide_data['blocks'].append({'type': 'bullet', 'content': '\n'.join(bullet_lines)})
                continue

            # Table block
            is_table = False
            if line.strip().startswith('|'):
                if (line_idx + 1 < len(lines)) and re.match(r'^\s*\|?.*--.*\|?\s*$', lines[line_idx+1]):
                    is_table = True
            
            if is_table:
                table_lines = []
                table_lines.append(lines[line_idx]) # Header
                line_idx += 1
                if line_idx < len(lines):
                    table_lines.append(lines[line_idx]) # Separator
                    line_idx += 1

                while line_idx < len(lines) and lines[line_idx].strip().startswith('|'):
                    table_lines.append(lines[line_idx]) # Data rows
                    line_idx += 1
                    
                slide_data['blocks'].append({'type': 'table', 'content': '\n'.join(table_lines)})
                continue

            # Regular text block
            text_lines = []
            while line_idx < len(lines):
                current_line = lines[line_idx]
                if not current_line.strip():
                    if line_idx + 1 < len(lines) and not lines[line_idx+1].lstrip().startswith(('-', '*', '+')) and not lines[line_idx+1].strip().startswith('|'):
                        text_lines.append(current_line)
                        line_idx += 1
                        continue
                    else:
                        break
                
                if current_line.lstrip().startswith(('-', '*', '+')): break
                if current_line.strip().startswith('|') and (line_idx + 1 < len(lines)) and re.match(r'^\s*\|?.*--.*\|?\s*$', lines[line_idx+1]): break
                    
                text_lines.append(current_line)
                line_idx += 1
            
            if text_lines:
                slide_data['blocks'].append({'type': 'text', 'content': '\n'.join(text_lines)})
            
            if not text_lines and not is_table and not (line.lstrip().startswith(('-', '*', '+'))):
                line_idx += 1 # Advance to prevent infinite loop
            
        slides.append(slide_data)
    return slides

def create_presentation_from_markdown(content: str, output_path: str = 'output.pptx') -> str:
    """Create a PowerPoint presentation from Markdown-structured text."""
    prs = Presentation()
    # Layout indices are theme-dependent, using common defaults
    layout_map = {
        'title_slide': prs.slide_layouts[0], 'title_content': prs.slide_layouts[1],
        'section_header': prs.slide_layouts[2], 'two_content': prs.slide_layouts[3],
        'comparison': prs.slide_layouts[4], 'title_only': prs.slide_layouts[5],
        'blank': prs.slide_layouts[6],
    }

    slides_data = parse_markdown_to_slides(content)

    for slide_data in slides_data:
        layout_name = slide_data['layout']
        slide_layout = layout_map.get(layout_name, layout_map['title_content'])
        current_slide = prs.slides.add_slide(slide_layout)

        if slide_data['title'] and current_slide.shapes.title:
            current_slide.shapes.title.text = slide_data['title']

        if layout_name in ['comparison', 'two_content']:
            content_placeholders = [p for p in current_slide.placeholders if p.placeholder_format.idx > 0 and p.has_text_frame]
            left_text_block = next((block for block in slide_data['blocks'] if block['type'] == 'text'), None)
            
            if left_text_block and len(content_placeholders) >= 2:
                left_ph, right_ph = content_placeholders[0], content_placeholders[1]
                text_block_content = left_text_block['content']
                
                left_text, right_text = text_block_content.split('|||', 1) if '|||' in text_block_content else (text_block_content, "")
                
                left_ph.text_frame.clear()
                right_ph.text_frame.clear()
                left_ph.text_frame.paragraphs[0].text = left_text.strip()
                right_ph.text_frame.paragraphs[0].text = right_text.strip()
            
        else:
            body_shape = next((shape for shape in current_slide.placeholders if shape.placeholder_format.idx != 0 and shape.has_text_frame), None)
            
            if body_shape:
                tf = body_shape.text_frame
                tf.clear()
                
                if len(tf.paragraphs) > 0 and not tf.paragraphs[0].text.strip():
                     tf.paragraphs[0].clear()

                for block in slide_data['blocks']:
                    if block['type'] == 'text':
                        p = tf.add_paragraph()
                        p.text = block['content']
                    elif block['type'] == 'bullet':
                        add_bullet_points_from_markdown(tf, block['content'])
                    elif block['type'] == 'table':
                        table_data = create_table_from_markdown(block['content'])
                        if not table_data: continue

                        rows, cols = len(table_data), len(table_data[0])
                        table_width = Inches(8.0)
                        table_height = Inches(0.4 * (rows + 1))
                        left = (prs.slide_width - table_width) / 2
                        top = Inches(2.0)
                        
                        try:
                            table_shape = current_slide.shapes.add_table(rows, cols, left, top, table_width, table_height)
                            table = table_shape.table
                            for r_idx, row_data in enumerate(table_data):
                                for c_idx, cell_text in enumerate(row_data):
                                    if c_idx < cols:
                                        table.cell(r_idx, c_idx).text = cell_text
                        except ValueError as e:
                            p = tf.add_paragraph()
                            p.text = f"[ERROR: Could not create table. Invalid data: {e}]"
                            
    prs.save(output_path)
    return output_path

if DOCLING_AVAILABLE:
    def extract_content_with_docling(file_path: str, page_range: str = None) :
        """Extracts content using the docling library."""
        try:
            ocr_options = TesseractCliOcrOptions(lang=["eng"])
            pipeline_options = PdfPipelineOptions(do_ocr=True, do_table_structure=True, ocr_options=ocr_options)
            doc_converter = DocumentConverter(
                format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=pipeline_options)})
            doc = doc_converter.convert(file_path).document
            
            if page_range:
                pages = []
                for part in page_range.split(','):
                    part = part.strip()
                    if '-' in part:
                        start, end = map(int, part.split('-'))
                        pages.extend(range(start - 1, end))
                    elif part.isdigit():
                        pages.append(int(part) - 1)
                
                filtered_content = []
                for i, page in enumerate(doc.pages):
                    if i in pages:
                        filtered_content.append(page.export_to_markdown())
                return "\n".join(filtered_content)
            else:
                return doc.export_to_markdown()

        except Exception as e:
            st.error(f"Error extracting content from {os.path.basename(file_path)}: {e}")
            return None


# --- Streamlit App Functions ---

def init_session_state():
    """Initialize Streamlit session state variables."""
    if 'resources' not in st.session_state:
        st.session_state['resources'] = []
    if 'markdown_content' not in st.session_state:
        st.session_state['markdown_content'] = ""
    # Add a key for the next resource
    if 'next_resource_key' not in st.session_state:
        st.session_state['next_resource_key'] = 0

def get_next_resource_key():
    key = st.session_state['next_resource_key']
    st.session_state['next_resource_key'] += 1
    return key

def add_file_resource(uploaded_file, page_range):
    """Saves uploaded file and adds it to resources."""
    if uploaded_file is not None:
        # Create a temporary file to save the uploaded content
        suffix = os.path.splitext(uploaded_file.name)[1] or '.tmp'
        temp_path = os.path.join(tempfile.gettempdir(), next(tempfile._get_candidate_names()) + suffix)

        try:
            # Write the UploadedFile object content to the temporary file
            file_bytes = uploaded_file.read()
            with open(temp_path, "wb") as f:
                f.write(file_bytes)
            
            st.session_state['resources'].append({
                "type": "File",
                "name": uploaded_file.name,
                "range": page_range,
                "temp_path": temp_path,
                "key": get_next_resource_key()
            })
            st.success(f"File '{uploaded_file.name}' added to resources.")
        except Exception as e:
            st.error(f"Failed to process uploaded file: {e}")
            if os.path.exists(temp_path): os.unlink(temp_path)

def add_url_resource(url, page_range):
    """Downloads content from URL, saves to temp, and adds to resources."""
    if not url:
        st.error("Please enter a URL.")
        return

    placeholder = st.empty()
    try:
        placeholder.info(f"Downloading content from {url}...")
        
        with httpx.Client(follow_redirects=True, timeout=15.0) as client:
            response = client.get(url)
            response.raise_for_status()

            # Create a temporary file to save the content
            suffix = os.path.splitext(url.split('?')[0])[1] or '.tmp'
            temp_path = os.path.join(tempfile.gettempdir(), next(tempfile._get_candidate_names()) + suffix)

            with open(temp_path, "wb") as temp_file:
                temp_file.write(response.content)

        st.session_state['resources'].append({
            "type": "URL",
            "name": url,
            "range": page_range,
            "temp_path": temp_path,
            "key": get_next_resource_key()
        })
        placeholder.success(f"URL '{url}' downloaded and added to resources.")
    except httpx.RequestError as e:
        placeholder.error(f"Could not fetch content from URL: {e}")
    except Exception as e:
        placeholder.error(f"An unexpected error occurred during URL processing: {e}")

def load_content_from_resources():
    """Extracts and concatenates content from all resources into the text area."""
    if not st.session_state['resources']:
        st.warning("No resources added to load content from.")
        return

    full_content = []
    
    with st.spinner("Extracting content... This may take a moment for large files or OCR."):
        for res in st.session_state['resources']:
            try:
                content = extract_content_with_docling(res["temp_path"], page_range=res["range"] or None)
                
                if content:
                    full_content.append(f"--- Resource: {res['name']} (Range: {res['range'] or 'All'}) ---\n{content}\n")
                else:
                    full_content.append(f"--- WARNING: No content extracted from {res['name']} ---\n")

            except Exception as e:
                full_content.append(f"--- ERROR processing {res['name']}: {e} ---\n")

    st.session_state['markdown_content'] = "\n\n".join(full_content)
    st.success("Content loaded successfully into the editor.")

def resource_manager_ui():
    """UI for adding, viewing, and managing resources."""
    
    st.subheader("1. Add Resources")
    
    # --- Add File ---
    with st.expander("Upload Document"):
        uploaded_file = st.file_uploader("Choose a file:", type=['pdf', 'docx', 'doc', 'txt', 'epub', 'html'], key="file_uploader")
        file_range = st.text_input("Page Range (e.g., 1-5, 2,4):", key="file_range_input")
        if st.button("Add File", key="add_file_btn"):
            if uploaded_file:
                add_file_resource(uploaded_file, file_range)
            else:
                st.warning("Please upload a file first.")

    # --- Add URL ---
    with st.expander("Import from URL"):
        url_input = st.text_input("Enter URL of a file:", key="url_input")
        url_range = st.text_input("Page Range (e.g., 1-5, 2,4):", key="url_range_input")
        if st.button("Add URL", key="add_url_btn"):
            add_url_resource(url_input, url_range)

    st.divider()
    st.subheader("2. Manage Resources")
    
    if st.session_state['resources']:
        
        # Prepare data for st.data_editor
        data_to_edit = [{
            'key': res['key'], 'Type': res['type'], 
            'Source': res['name'], 'Range/Chapter': res['range']
        } for res in st.session_state['resources']]
        
        edited_data = st.data_editor(
            data_to_edit,
            column_order=('Type', 'Source', 'Range/Chapter'),
            column_config={"key": st.column_config.Column("Key", disabled=True, width="small"),
                           "Type": st.column_config.Column("Type", disabled=True, width="small"),
                           "Source": st.column_config.Column("Source", disabled=True, width="large"),
                           "Range/Chapter": st.column_config.TextColumn("Range/Chapter", width="medium")},
            hide_index=True,
            num_rows="dynamic",
            key="resource_data_editor"
        )
        
        # Check for deleted rows and update ranges
        edited_keys = {item['key'] for item in edited_data}
        
        # 1. Update ranges for existing resources
        for edited_row in edited_data:
            original_resource = next((res for res in st.session_state['resources'] if res['key'] == edited_row['key']), None)
            if original_resource:
                original_resource['range'] = edited_row['Range/Chapter']

        # 2. Check for deletions
        if len(edited_data) < len(st.session_state['resources']):
            deleted_resources = [res for res in st.session_state['resources'] if res['key'] not in edited_keys]
            
            # Clean up temp files for deleted resources
            for res in deleted_resources:
                if 'temp_path' in res and os.path.exists(res['temp_path']):
                    os.unlink(res['temp_path'])
            
            st.session_state['resources'] = [res for res in st.session_state['resources'] if res['key'] in edited_keys]
            st.success("Resource(s) removed.")
        
        st.divider()
        if st.button("Load Content from Resources (Step 3)", type="primary"):
            load_content_from_resources()

    else:
        st.info("No resources added yet.")


def main_app_ui():
    """Main UI for content editing and presentation generation."""
    
    st.title("📄 Presentation Generator")
    st.caption("Upload documents or import from a URL, refine the content in Markdown, and generate a PowerPoint presentation.")
    
    # Resources in the Sidebar
    with st.sidebar:
        resource_manager_ui()
    
    # Main Content Area
    st.header("3. Presentation Content (Editable Markdown)")
    st.caption("Content is loaded here from resources. Use the custom Markdown syntax (layouts, tables, bullets, --- for slide breaks) as shown in the example in `pptx_utils.py`.")
    
    # Text Area
    markdown_content = st.text_area(
        "Edit Content:", 
        st.session_state['markdown_content'],
        height=500,
        key="markdown_editor"
    )
    # Important: Update session state when text area is changed
    st.session_state['markdown_content'] = markdown_content

    st.header("4. Generate Presentation")
    
    # Generate Button
    if st.button("Generate Presentation (.pptx)", type="success", key="generate_btn"):
        content = st.session_state['markdown_content']
        if not content.strip():
            st.error("The content editor is empty.")
            return
        
        temp_output_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
        output_path = temp_output_file.name
        temp_output_file.close()
        
        try:
            with st.spinner("Creating PowerPoint..."):
                create_presentation_from_markdown(content, output_path)
            
            with open(output_path, "rb") as f:
                pptx_bytes = f.read()

            st.download_button(
                label="✅ Download Presentation",
                data=pptx_bytes,
                file_name="generated_presentation.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="download_button",
                type="primary"
            )
            st.balloons()
            st.success("Presentation created successfully! Click the download button above.")
            
        except Exception as e:
            st.error(f"Failed to generate presentation: {e}")
            st.exception(e)
        finally:
            if os.path.exists(output_path):
                os.unlink(output_path)

# --- Run App ---
if __name__ == "__main__":
    init_session_state()
    main_app_ui()
